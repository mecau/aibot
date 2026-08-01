import asyncio
import logging
import base64
import io
import json
import os
import re
from aiogram import Bot, Dispatcher, types, F
from aiogram.filters import CommandStart, Command
from aiogram.types import (
    InlineKeyboardMarkup, InlineKeyboardButton,
    ReplyKeyboardMarkup, KeyboardButton, Message, BotCommand, BufferedInputFile
)
from openai import AsyncOpenAI
from docx import Document
from docx.shared import Pt, Inches, Mm
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

from config import (
    BOT_TOKEN, GROQ_API_KEY, CHANNEL_USERNAME, CHANNEL_URL,
    TEXT_MODEL, VISION_MODEL, PROMPTS, AD_FOOTER
)

MY_ADMIN_ID = 1184589026

logging.basicConfig(level=logging.INFO)
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher()

groq_client = AsyncOpenAI(
    api_key=GROQ_API_KEY,
    base_url="https://api.groq.com/openai/v1"
)

USERS_FILE = "users.json"
FAV_FILE = "favorites.json"

def load_user_ids() -> set:
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, "r") as f:
                return set(json.load(f))
        except Exception:
            pass
    return set()

def save_user_id(user_id: int):
    ids = load_user_ids()
    if user_id not in ids:
        ids.add(user_id)
        with open(USERS_FILE, "w") as f:
            json.dump(list(ids), f)

def load_favorites() -> dict:
    if os.path.exists(FAV_FILE):
        try:
            with open(FAV_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {}

def save_favorites(favs: dict):
    with open(FAV_FILE, "w", encoding="utf-8") as f:
        json.dump(favs, f, ensure_ascii=False, indent=2)

def add_to_favorites(user_id: int, text: str):
    favs = load_favorites()
    uid_str = str(user_id)
    if uid_str not in favs:
        favs[uid_str] = []
    if text not in favs[uid_str]:
        favs[uid_str].append(text)
        save_favorites(favs)

users_db = {}
broadcast_states = set()
ppt_states = set()
excel_states = set()

def get_user_data(user_id: int):
    if user_id not in users_db:
        users_db[user_id] = {
            "mode": "ai",
            "history": [],
            "last_output": "Здесь пока нет ответов."
        }
    return users_db[user_id]

async def check_subscription(user_id: int) -> bool:
    try:
        member = await bot.get_chat_member(chat_id=f"@{CHANNEL_USERNAME}", user_id=user_id)
        if member.status in ["creator", "administrator", "member"]:
            return True
    except Exception:
        pass
    return False

def get_sub_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📢 Подписаться на канал", url=CHANNEL_URL)],
        [InlineKeyboardButton(text="✅ Проверить подписку", callback_data="check_sub")]
    ])

def get_main_keyboard():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="📄 Скачать ответ в Word"), KeyboardButton(text="📑 Создать титульник ГОСТ")],
            [KeyboardButton(text="📈 Создать презентацию"), KeyboardButton(text="📊 Создать таблицу Excel")],
            [KeyboardButton(text="⭐ Избранное"), KeyboardButton(text="🔄 Сменить режим")],
            [KeyboardButton(text="ℹ️ О MecauAI")]
        ],
        resize_keyboard=True
    )

def get_admin_keyboard():
    return ReplyKeyboardMarkup(
        keyboard=[
            [KeyboardButton(text="📄 Скачать ответ в Word"), KeyboardButton(text="📑 Создать титульник ГОСТ")],
            [KeyboardButton(text="📈 Создать презентацию"), KeyboardButton(text="📊 Создать таблицу Excel")],
            [KeyboardButton(text="⭐ Избранное"), KeyboardButton(text="🔄 Сменить режим")],
            [KeyboardButton(text="ℹ️ О MecauAI")],
            [KeyboardButton(text="📊 Статистика бота"), KeyboardButton(text="📢 Сделать рассылку")]
        ],
        resize_keyboard=True
    )

def keyboard_for(user_id: int):
    return get_admin_keyboard() if user_id == MY_ADMIN_ID else get_main_keyboard()

def get_answer_inline_keyboard():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="💡 Объяснить проще", callback_data="btn_simplify"),
            InlineKeyboardButton(text="⭐ В избранное", callback_data="btn_save_fav")
        ]
    ])

def clean_text(text: str) -> str:
    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
    text = text.replace("**", "").replace("*", "")
    return text.strip()

@dp.message(CommandStart())
async def cmd_start(message: Message):
    user_id = message.from_user.id
    save_user_id(user_id)

    if user_id != MY_ADMIN_ID:
        if not await check_subscription(user_id):
            await message.answer(
                f"🔒 Доступ заблокирован!\n\n"
                f"Чтобы пользоваться MecauAI, подпишись на наш канал:\n"
                f" 👉 {CHANNEL_URL}\n\n"
                f"После подписки нажми кнопку ниже 👇",
                reply_markup=get_sub_keyboard()
            )
            return

    start_text = (
        f"Привет, {message.from_user.first_name}! Ты активировал MecauAI 🚀\n\n"
        "Я твой карманный помощник для учебы. Отправляй любые вопросы, задачи, конспекты, а также создавай презентации и таблицы!"
    )
    await message.answer(start_text, reply_markup=keyboard_for(user_id))

@dp.callback_query(F.data == "check_sub")
async def cb_check_sub(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    if await check_subscription(user_id):
        save_user_id(user_id)
        await callback.message.delete()
        await callback.message.answer(
            "🎉 Подписка подтверждена! Добро пожаловать в MecauAI 🚀",
            reply_markup=keyboard_for(user_id)
        )
    else:
        await callback.answer("❌ Ты еще не подписался на канал!", show_alert=True)

@dp.message(F.text == "ℹ️ О MecauAI")
@dp.message(Command("about"))
async def cmd_about(message: Message):
    if message.from_user.id != MY_ADMIN_ID and not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return
    about_text = (
        "🧠 Возможности MecauAI:\n\n"
        "• AI-Ассистент и Друг: Помогает учиться и поддерживает в дедлайны.\n"
        "• Анализ фото/файлов/текста: Быстрый разбор материалов и конспектов.\n"
        "• 📄 Экспорт в .docx: Сохранение любого ответа в Word.\n"
        "• 📑 Титульник по ГОСТу: Быстрое оформление титульных листов.\n"
        "• 📈 Генерация презентаций (.pptx): Автоматическое создание слайдов.\n"
        "• 📊 Создание таблиц Excel (.xlsx): Аккуратные структурированные данные.\n"
        "• ⭐ Избранное: Сохранение важных ответов под рукой."
    )
    await message.answer(about_text, reply_markup=keyboard_for(message.from_user.id))

@dp.message(F.text == "🔄 Сменить режим")
@dp.message(Command("mode"))
async def cmd_mode(message: Message):
    if message.from_user.id != MY_ADMIN_ID and not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🧠 Академический ассистент", callback_data="set_mode_ai")],
        [InlineKeyboardButton(text="🫂 Лучший друг",             callback_data="set_mode_friend")]
    ])
    await message.answer("Выбери режим работы бота:", reply_markup=kb)

@dp.callback_query(F.data.startswith("set_mode_"))
async def cb_set_mode(callback: types.CallbackQuery):
    mode = callback.data.split("_")[-1]
    user_data = get_user_data(callback.from_user.id)
    user_data["mode"] = mode
    name = "Академический ассистент 🧠" if mode == "ai" else "Лучший друг 🫂"
    await callback.message.edit_text(f"Режим переключен на: {name}")
    await callback.answer()

@dp.message(F.text == "📊 Статистика бота")
async def cmd_stats(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    users = load_user_ids()
    await message.answer(f"📊 Статистика бота:\n\n👥 Всего пользователей: {len(users)}")

@dp.message(F.text == "📢 Сделать рассылку")
async def cmd_broadcast_prompt(message: Message):
    if message.from_user.id != MY_ADMIN_ID:
        return
    broadcast_states.add(message.from_user.id)
    await message.answer("📢 Отправь текст рассылки следующим сообщением (все пользователи бота получат его).")

@dp.message(F.text == "📈 Создать презентацию")
async def cmd_ppt_prompt(message: Message):
    user_id = message.from_user.id
    if user_id != MY_ADMIN_ID and not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return
    ppt_states.add(user_id)
    await message.answer("📈 Отправь тему или текст для презентации следующим сообщением, и я соберу слайды (.pptx)!")

@dp.message(F.text == "📊 Создать таблицу Excel")
async def cmd_excel_prompt(message: Message):
    user_id = message.from_user.id
    if user_id != MY_ADMIN_ID and not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return
    excel_states.add(user_id)
    await message.answer("📊 Опиши задачу или тему для таблицы следующим сообщением, и я сформирую аккуратный Excel-файл (.xlsx)!")

@dp.message(F.text == "⭐ Избранное")
async def cmd_favorites(message: Message):
    user_id = message.from_user.id
    if user_id != MY_ADMIN_ID and not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return

    favs = load_favorites()
    uid_str = str(user_id)
    user_favs = favs.get(uid_str, [])

    if not user_favs:
        await message.answer("⭐ У тебя пока нет сохраненных ответов в избранном.\nНажимай кнопку «⭐ В избранное» под сообщениями бота!")
        return

    await message.answer(f"⭐ Твои сохраненные ответы ({len(user_favs)}):")
    for idx, fav_text in enumerate(user_favs, 1):
        display_text = f"Сохранение #{idx}:\n\n{fav_text}"
        if len(display_text) > 4000:
            display_text = display_text[:3997] + "..."
        await message.answer(display_text, disable_web_page_preview=True)

@dp.callback_query(F.data.in_({"btn_simplify", "btn_save_fav"}))
async def cb_answer_actions(callback: types.CallbackQuery):
    user_id = callback.from_user.id
    user_data = get_user_data(user_id)
    
    msg_text = callback.message.text or callback.message.caption or ""
    
    if "—\n⚡" in msg_text:
        clean_msg = msg_text.split("—\n⚡")[0].strip()
    else:
        clean_msg = msg_text.strip()

    if not clean_msg:
        await callback.answer("⚠️ Нечего обрабатывать!", show_alert=True)
        return

    if callback.data == "btn_save_fav":
        add_to_favorites(user_id, clean_msg)
        await callback.answer("⭐ Успешно сохранено в избранное!", show_alert=True)

    elif callback.data == "btn_simplify":
        await callback.answer("💡 Сжимаю до сути...")
        try:
            response = await groq_client.chat.completions.create(
                model=TEXT_MODEL,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "Ты — помощник, который объясняет сложные вещи предельно просто и коротко. "
                            "ПРАВИЛО: дай краткое объяснение в 3–5 предложениях максимум — как другу на пальцах. "
                            "Никаких длинных лекций, никакой воды. Только суть, простым языком."
                        )
                    },
                    {"role": "user", "content": f"Объясни коротко и просто, самую суть:\n\n{clean_msg}"}
                ],
                temperature=0.7
            )
            simplified_reply = clean_text(response.choices[0].message.content)
            
            user_data["last_output"] = simplified_reply
            
            if user_data["history"]:
                user_data["history"][-1]["content"] = f"[Упрощено до сути]: {simplified_reply}"

            full_reply = f"💡 <b>Коротко на пальцах:</b>\n\n{simplified_reply}{AD_FOOTER}"
            await callback.message.answer(
                full_reply,
                reply_markup=get_answer_inline_keyboard(),
                parse_mode="HTML",
                disable_web_page_preview=True
            )
        except Exception as e:
            await callback.message.answer(f"⚠️ Ошибка при упрощении: {e}")

@dp.message(F.text == "📄 Скачать ответ в Word")
async def cmd_download_word(message: Message):
    user_id = message.from_user.id
    if user_id != MY_ADMIN_ID and not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return

    user_data = get_user_data(user_id)
    text_to_save = user_data.get("last_output", "").strip()

    if not text_to_save or text_to_save == "Здесь пока нет ответов.":
        await message.answer("⚠️ У тебя еще нет ответов от ИИ, которые можно сохранить. Сначала отправь запрос!")
        return

    doc = Document()
    for section in doc.sections:
        section.top_margin = Mm(20)
        section.bottom_margin = Mm(20)
        section.left_margin = Mm(30)
        section.right_margin = Mm(15)

    p = doc.add_paragraph()
    run = p.add_run(text_to_save)
    run.font.name = 'Times New Roman'
    run.font.size = Pt(14)

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)

    file_doc = BufferedInputFile(bio.read(), filename="MecauAI_Answer.docx")
    await message.answer_document(file_doc, caption="📄 Вот твой ответ в формате Word!")

@dp.message(F.text == "📑 Создать титульник ГОСТ")
async def cmd_gost_title(message: Message):
    if message.from_user.id != MY_ADMIN_ID and not await check_subscription(message.from_user.id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📘 Индивидуальный проект", callback_data="gost_project")],
        [InlineKeyboardButton(text="📗 Курсовая работа",        callback_data="gost_coursework")],
        [InlineKeyboardButton(text="📕 Дипломная работа (ВКР)", callback_data="gost_diploma")],
        [InlineKeyboardButton(text="📙 Отчёт по практике",      callback_data="gost_practice")],
    ])
    await message.answer("📑 Выбери тип работы для титульника:", reply_markup=kb)

@dp.callback_query(F.data.startswith("gost_"))
async def cb_gost_generate(callback: types.CallbackQuery):
    work_type_map = {
        "gost_project":    ("ИНДИВИДУАЛЬНЫЙ ПРОЕКТ",                   "Индивидуальный_проект"),
        "gost_coursework": ("КУРСОВАЯ РАБОТА",                         "Курсовая_работа"),
        "gost_diploma":    ("ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА (ВКР)", "ВКР"),
        "gost_practice":   ("ОТЧЁТ ПО ПРАКТИКЕ",                      "Отчет_по_практике"),
    }
    work_label, filename_base = work_type_map[callback.data]

    doc = Document()
    for section in doc.sections:
        section.top_margin    = Mm(20)
        section.bottom_margin = Mm(20)
        section.left_margin   = Mm(30)
        section.right_margin  = Mm(15)

    def add_centered(text, size=14, bold=False):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.name = 'Times New Roman'
        run.font.size = Pt(size)
        run.bold = bold
        return p

    def add_right(text, size=14):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = p.add_run(text)
        run.font.name = 'Times New Roman'
        run.font.size = Pt(size)
        return p

    def add_empty(count=1):
        for _ in range(count):
            p = doc.add_paragraph()
            p.add_run("").font.name = 'Times New Roman'

    add_centered(
        "ФЕДЕРАЛЬНОЕ ГОСУДАРСТВЕННОЕ БЮДЖЕТНОЕ ПРОФЕССИОНАЛЬНОЕ\n"
        "ОБРАЗОВАТЕЛЬНОЕ УЧРЕЖДЕНИЕ\n"
        "«НАЗВАНИЕ КОЛЛЕДЖА»",
        size=14
    )
    add_empty(4)
    add_centered(work_label, size=14, bold=True)
    add_empty(1)
    add_centered("на тему:\n«Введи тему работы здесь»", size=14, bold=True)
    add_empty(6)
    add_right(
        "Выполнил(а): студент(ка) группы ГРУППА\n"
        "Фамилия Имя Отчество\n\n"
        "Руководитель:\n"
        "Должность, Фамилия И.О."
    )
    add_empty(5)
    add_centered("Город — 2026", size=14)

    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)

    file_doc = BufferedInputFile(bio.read(), filename=f"Titulnik_{filename_base}.docx")
    await callback.message.answer_document(file_doc, caption=f"📑 Титульник «{work_label}» готов!")
    await callback.answer()

MENU_BUTTONS = {
    "📄 Скачать ответ в Word", "📑 Создать титульник ГОСТ",
    "📈 Создать презентацию", "📊 Создать таблицу Excel",
    "⭐ Избранное", "🔄 Сменить режим", "ℹ️ О MecauAI", 
    "📢 Сделать рассылку", "📊 Статистика бота"
}

@dp.message(F.document)
async def handle_document(message: Message):
    user_id = message.from_user.id
    if user_id != MY_ADMIN_ID and not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return

    save_user_id(user_id)
    doc_file = message.document
    file_name = doc_file.file_name or "file"
    
    if doc_file.file_size and doc_file.file_size > 10 * 1024 * 1024:
        await message.answer("⚠️ Файл слишком большой. Максимальный размер — 10 МБ.")
        return

    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        file_info = await bot.get_file(doc_file.file_id)
        downloaded = await bot.download_file(file_info.file_path)
        file_bytes = downloaded.read()

        file_content_text = ""
        
        if file_name.endswith(".docx"):
            doc_obj = Document(io.BytesIO(file_bytes))
            file_content_text = "\n".join([p.text for p in doc_obj.paragraphs if p.text.strip()])
        else:
            try:
                file_content_text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                file_content_text = file_bytes.decode("cp1251", errors="ignore")

        if not file_content_text.strip():
            await message.answer("⚠️ Не удалось прочитать текст из этого файла или он пустой.")
            return

        if len(file_content_text) > 15000:
            file_content_text = file_content_text[:15000] + "\n[...текст файла обрезан из-за большой длины...]"

        user_prompt = message.caption or "Проанализируй содержимое этого файла, сделай выводы или реши задачу в нем."
        
        user_data = get_user_data(user_id)
        system_prompt = PROMPTS[user_data["mode"]]

        prompt_payload = (
            f"Пользователь прикрепил файл '{file_name}'. Вот его содержимое:\n\n"
            f"```\n{file_content_text}\n```\n\n"
            f"Задание пользователя к файлу: {user_prompt}"
        )

        user_data["history"].append({"role": "user", "content": prompt_payload})
        if len(user_data["history"]) > 6:
            user_data["history"] = user_data["history"][-6:]

        strict_system_prompt = (
            system_prompt + 
            "\n\nВНИМАНИЕ: Категорически запрещено отвечать списком в 2-3 слова или телеграфным стилем. "
            "Каждый ответ должен быть развернутой статьей с полными предложениями, абзацами и структурой."
        )
        messages_payload = [{"role": "system", "content": strict_system_prompt}] + user_data["history"]

        response = await groq_client.chat.completions.create(
            model=TEXT_MODEL,
            messages=messages_payload,
            temperature=0.7
        )
        ai_reply = clean_text(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply

        full_message = f"{ai_reply}{AD_FOOTER}"

        await message.answer(full_message, parse_mode="HTML", reply_markup=get_answer_inline_keyboard(), disable_web_page_preview=True)
    except Exception as e:
        await message.answer(f"⚠️ Ошибка при обработке файла: {e}")

@dp.message(F.photo)
async def handle_photo(message: Message):
    user_id = message.from_user.id
    if user_id != MY_ADMIN_ID and not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return

    save_user_id(user_id)
    prompt = message.caption or "Проанализируй это учебное изображение, сделай конспект и выдели суть."
    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        photo = message.photo[-1]
        file_info = await bot.get_file(photo.file_id)
        downloaded_file = await bot.download_file(file_info.file_path)
        base64_image = base64.b64encode(downloaded_file.read()).decode('utf-8')
        
        response = await groq_client.chat.completions.create(
            model=VISION_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                ]
            }]
        )
        reply = clean_text(response.choices[0].message.content)
        
        user_data = get_user_data(user_id)
        user_data["last_output"] = reply

        full_message = f"{reply}{AD_FOOTER}"

        await message.answer(full_message, parse_mode="HTML", reply_markup=get_answer_inline_keyboard(), disable_web_page_preview=True)
    except Exception as e:
        await message.answer(f"⚠️ Ошибка анализа фото: {e}")

@dp.message(F.text)
async def handle_text(message: Message):
    user_id = message.from_user.id

    if user_id == MY_ADMIN_ID and user_id in broadcast_states:
        broadcast_states.remove(user_id)
        users = load_user_ids()
        success, failed = 0, 0
        
        status_msg = await message.answer("📢 Начинаю рассылку...")
        for uid in users:
            try:
                await bot.send_message(uid, message.text, disable_web_page_preview=True)
                success += 1
                await asyncio.sleep(0.05)
            except Exception:
                failed += 1
        
        await status_msg.edit_text(f"✅ Рассылка завершена!\n\n👥 Доставлено: {success}\n❌ Ошибок: {failed}")
        return

    if user_id in ppt_states:
        ppt_states.remove(user_id)
        await bot.send_chat_action(chat_id=message.chat.id, action="typing")
        status_msg = await message.answer("📈 Генерирую презентацию, пишу развернутый текст...")
        try:
            prompt = message.text
            response = await groq_client.chat.completions.create(
                model=TEXT_MODEL,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "Ты — профессиональный академический методист. На основе темы пользователя "
                            "составь подробную презентацию из 6-8 слайдов. "
                            "ВНИМАНИЕ: Текст на слайдах (points) должен быть развернутым и содержательным! "
                            "Пиши полными предложениями (от 15 до 30 слов в каждом пункте), раскрывай суть, чтобы не было пустых слайдов. "
                            "Ответ выдай строго в формате JSON без лишнего текста: "
                            "[{\"title\": \"Заголовок\", \"points\": [\"Развернутый тезис 1...\", \"Развернутый тезис 2...\"]}]"
                        )
                    },
                    {"role": "user", "content": f"Тема презентации: {prompt}"}
                ],
                temperature=0.7
            )
            raw_content = clean_text(response.choices[0].message.content)
            if "```json" in raw_content:
                raw_content = raw_content.split("```json")[1].split("```")[0].strip()
            elif "```" in raw_content:
                raw_content = raw_content.split("```")[1].split("```")[0].strip()
            
            slides_data = json.loads(raw_content)

            prs = Presentation()
            slide_layout = prs.slide_layouts[1]

            title_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = "Презентация"
            if title_shape.text_frame.paragraphs:
                title_shape.text_frame.paragraphs[0].font.name = 'Times New Roman'
                title_shape.text_frame.paragraphs[0].font.size = PptxPt(40)

            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = prompt
            if subtitle_shape.text_frame.paragraphs:
                subtitle_shape.text_frame.paragraphs[0].font.name = 'Times New Roman'
                subtitle_shape.text_frame.paragraphs[0].font.size = PptxPt(24)

            for item in slides_data:
                slide = prs.slides.add_slide(slide_layout)
                
                title_shape = slide.shapes.title
                title_shape.text = item.get("title", "Слайд")
                if title_shape.text_frame.paragraphs:
                    title_shape.text_frame.paragraphs[0].font.name = 'Times New Roman'
                    title_shape.text_frame.paragraphs[0].font.size = PptxPt(32)
                
                tf = slide.placeholders[1].text_frame
                tf.text = ""
                for pt in item.get("points", []):
                    p = tf.add_paragraph()
                    p.text = pt
                    p.level = 0
                    p.font.name = 'Times New Roman'
                    p.font.size = PptxPt(20)

            bio = io.BytesIO()
            prs.save(bio)
            bio.seek(0)

            file_doc = BufferedInputFile(bio.read(), filename="Presentation.pptx")
            await status_msg.delete()
            await message.answer_document(file_doc, caption="📈 Твоя развернутая презентация готова!")
        except Exception as e:
            await status_msg.delete()
            await message.answer(f"⚠️ Не удалось сгенерировать презентацию: {e}")
        return

    if user_id in excel_states:
        excel_states.remove(user_id)
        await bot.send_chat_action(chat_id=message.chat.id, action="typing")
        status_msg = await message.answer("📊 Создаю структуру таблицы Excel...")
        try:
            prompt = message.text
            response = await groq_client.chat.completions.create(
                model=TEXT_MODEL,
                messages=[
                    {
                        "role": "system",
                        "content": (
                            "Ты аналитик данных. На основе запроса пользователя создай подробную и реалистичную таблицу Excel. "
                            "Она должна содержать как минимум 4-5 столбцов и 5-10 строк данных (если это уместно). "
                            "Ответ дай строго в формате JSON со списком массивов (строк), где первая строка — заголовки столбцов. "
                            "Пример формата: [[\"Колонка 1\", \"Колонка 2\"], [\"Данные 1\", \"Данные 2\"]. "
                            "Никакого другого текста, только JSON."
                        )
                    },
                    {"role": "user", "content": f"Запрос для таблицы: {prompt}"}
                ],
                temperature=0.7
            )
            raw_content = clean_text(response.choices[0].message.content)
            if "```json" in raw_content:
                raw_content = raw_content.split("```json")[1].split("```")[0].strip()
            elif "```" in raw_content:
                raw_content = raw_content.split("```")[1].split("```")[0].strip()

            table_data = json.loads(raw_content)

            wb = Workbook()
            ws = wb.active
            ws.title = "Данные"

            header_fill = PatternFill(start_color="1F4E78", end_color="1F4E78", fill_type="solid")
            header_font = Font(name="Times New Roman", size=12, bold=True, color="FFFFFF")
            regular_font = Font(name="Times New Roman", size=12)
            thin_border = Border(
                left=Side(style='thin', color='D9D9D9'),
                right=Side(style='thin', color='D9D9D9'),
                top=Side(style='thin', color='D9D9D9'),
                bottom=Side(style='thin', color='D9D9D9')
            )

            for row_idx, row in enumerate(table_data, 1):
                ws.append(row)
                for col_idx in range(1, len(row) + 1):
                    cell = ws.cell(row=row_idx, column=col_idx)
                    cell.border = thin_border
                    if row_idx == 1:
                        cell.fill = header_fill
                        cell.font = header_font
                        cell.alignment = Alignment(horizontal="center", vertical="center")
                    else:
                        cell.font = regular_font
                        cell.alignment = Alignment(horizontal="left", vertical="center")

            for col in ws.columns:
                max_len = max(len(str(cell.value or '')) for cell in col)
                col_letter = col[0].column_letter
                ws.column_dimensions[col_letter].width = max(max_len + 2, 12)

            bio = io.BytesIO()
            wb.save(bio)
            bio.seek(0)

            file_doc = BufferedInputFile(bio.read(), filename="Table.xlsx")
            await status_msg.delete()
            await message.answer_document(file_doc, caption="📊 Твоя таблица Excel готова!")
        except Exception as e:
            await status_msg.delete()
            await message.answer(f"⚠️ Не удалось создать таблицу: {e}")
        return

    if message.text in MENU_BUTTONS:
        return

    if user_id != MY_ADMIN_ID and not await check_subscription(user_id):
        await message.answer(f"🔒 Сначала подпишись на канал:\n{CHANNEL_URL}", reply_markup=get_sub_keyboard())
        return

    save_user_id(user_id)
    user_data = get_user_data(user_id)
    system_prompt = PROMPTS[user_data["mode"]]

    user_data["history"].append({"role": "user", "content": message.text})
    if len(user_data["history"]) > 6:
        user_data["history"] = user_data["history"][-6:]

    strict_system_prompt = (
        system_prompt + 
        "\n\nВНИМАНИЕ: Категорически запрещено отвечать списком в 2-3 слова или телеграфным стилем. "
        "Каждый ответ должен быть развернутой статьей с полными предложениями, абзацами и структурой."
    )
    messages_payload = [{"role": "system", "content": strict_system_prompt}] + user_data["history"]

    await bot.send_chat_action(chat_id=message.chat.id, action="typing")
    try:
        response = await groq_client.chat.completions.create(
            model=TEXT_MODEL,
            messages=messages_payload,
            temperature=0.7
        )
        ai_reply = clean_text(response.choices[0].message.content)
        if not ai_reply:
            ai_reply = "Готово."

        user_data["history"].append({"role": "assistant", "content": ai_reply})
        user_data["last_output"] = ai_reply

        full_message = f"{ai_reply}{AD_FOOTER}"

        await message.answer(full_message, parse_mode="HTML", reply_markup=get_answer_inline_keyboard(), disable_web_page_preview=True)
    except Exception as e:
        await message.answer(f"⚠️ Ошибка обработки: {e}")

async def main():
    await bot.set_my_commands([
        BotCommand(command="start", description="Перезапустить бота"),
        BotCommand(command="mode",  description="Сменить режим ассистента"),
        BotCommand(command="about", description="О возможностях"),
    ])
    
    await bot.delete_webhook(drop_pending_updates=True)
    
    # Бесконечный цикл с автоперезапуском при обрывах связи
    while True:
        try:
            print("🚀 Бот MecauAI запущен и слушает обновления...")
            await dp.start_polling(bot)
        except Exception as e:
            print(f"⚠️ Ошибка соединения: {e}. Перезапуск через 5 секунд...")
            await asyncio.sleep(5)

if __name__ == "__main__":
    asyncio.run(main())
